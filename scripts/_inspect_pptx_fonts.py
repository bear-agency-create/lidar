# -*- coding: utf-8 -*-
"""Inspect fonts/colors and dump full slide texts from template PPTX."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

pptx = Path(r"c:\Users\user\Downloads\Telegram Desktop\AeroPorter_презентация (2).pptx")
prs = Presentation(str(pptx))

def color_of(run):
    try:
        c = run.font.color
        if c is None or c.rgb is None:
            return None
        return str(c.rgb)
    except Exception:
        return None

for i, slide in enumerate(prs.slides, 1):
    print("=" * 70)
    print("SLIDE", i)
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for pi, p in enumerate(shape.text_frame.paragraphs):
            raw = "".join(r.text for r in p.runs) if p.runs else (p.text or "")
            if not raw.strip():
                continue
            sizes, colors, names, bolds = [], [], [], []
            for r in p.runs:
                sizes.append(r.font.size.pt if r.font.size else None)
                colors.append(color_of(r))
                names.append(r.font.name)
                bolds.append(r.font.bold)
            print(f"  [{shape.name}] {raw!r}")
            print(f"      size={sizes} bold={bolds} font={names} color={colors}")
