# -*- coding: utf-8 -*-
"""Build English AeroPorter/AirPorter deck from the RU design template,
aligned to the Final Speech (RU + EN) content."""
from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation

SRC = Path(r"c:\Users\user\Downloads\Telegram Desktop\AeroPorter_презентация (2).pptx")
OUT = Path(r"c:\Users\user\Desktop\AirPorter — Presentation EN.pptx")
DOCX_SRC = Path(r"c:\Users\user\Downloads\Telegram Desktop\AirPorter — Final Speech (RU + EN).docx")
DOCX_OUT = Path(r"c:\Users\user\Desktop\AirPorter — Final Speech (RU + EN).docx")

# Exact RU strings from the template → EN aligned to the speech.
REPLACEMENTS: list[tuple[str, str]] = [
    # Slide 1 — cover
    ("РОБОТОТЕХНИКА · ПРОЕКТНАЯ РАБОТА", "ROBOTICS · PROJECT WORK"),
    ("AeroPorter", "AirPorter"),
    ("Мобильный робот-ассистент", "A mobile robot assistant"),
    ("для аэропорта", "for the airport"),
    (
        "Цель: считать билет · показать рейс · построить маршрут сопровождения",
        "Goal: scan a ticket · show the flight · build an escort route",
    ),
    ("Команда", "Team"),
    ("Илья Карякин · 9 класс", "Ilya Karyakin · Grade 9"),
    ("Парамонов Станислав · 9 класс", "Stanislav Paramonov · Grade 9"),
    ("Сабирзанов Артур · 11 класс", "Artur Sabirzanov · Grade 11"),
    ("Руководитель", "Supervisor"),
    ("Александр Борисович", "Alexander Borisovich"),
    ("Четвергов", "Chetvergov"),
    ("AEROPORTER", "AIRPORTER"),

    # Slide 2 — problem
    ("01 · ПРОБЛЕМА", "01 · THE PROBLEM"),
    ("Аэропорт сложен для пассажира", "An airport is hard for a passenger"),
    (
        "Информация меняется быстро, а помощь нужна здесь и сейчас",
        "Information changes fast, and help is needed here and now",
    ),
    ("Навигация", "Navigation"),
    (
        "Поиск стойки, выхода и багажа отнимает время.",
        "Finding check-in, the gate, and baggage takes time.",
    ),
    ("Информация", "Information"),
    (
        "Статус рейса и номер выхода могут измениться.",
        "Flight status and gate numbers can change.",
    ),
    ("Доступность", "Accessibility"),
    (
        "Язык, возраст и стресс усложняют общение.",
        "Language, age, and stress make communication harder.",
    ),
    ("Особенно уязвимы", "Especially affected"),
    ("пожилые люди", "elderly people"),
    ("семьи с детьми", "families with children"),
    ("иностранные пассажиры", "international passengers"),

    # Slide 3 — solution
    ("02 · РЕШЕНИЕ", "02 · THE SOLUTION"),
    (
        "Один понятный путь: от билета до нужного места",
        "One clear path: from ticket to destination",
    ),
    ("Сканирование", "Scanning"),
    ("Камера считывает штрихкод билета", "The camera reads the ticket barcode"),
    ("Экран показывает рейс, время и выход", "The screen shows flight, time, and gate"),
    ("Маршрут", "Route"),
    (
        "Билет связывается с настроенной точкой",
        "The ticket is linked to a configured point",
    ),
    ("Сопровождение", "Escort"),
    ("Робот запускает маршрут A→B", "The robot starts route A→B"),
    (
        "Цель: снизить стресс пассажира и нагрузку на персонал",
        "Goal: reduce passenger stress and staff workload",
    ),

    # Slide 4 — hardware
    ("03 · ПРОТОТИП", "03 · PROTOTYPE"),
    ("Аппаратная платформа", "Hardware platform"),
    (
        "Манёвренное шасси и два уровня управления",
        "A maneuverable chassis and two control layers",
    ),
    ("обработка изображения и интерфейс", "image processing and interface"),
    ("управление моторами и watchdog команд", "motor control and command watchdog"),
    ("4 Mecanum-колеса", "4 Mecanum wheels"),
    (
        "движение вперёд, вбок и с разворотом",
        "forward, sideways, and in-place turns",
    ),
    ("Привод и датчики", "Drive and sensors"),
    (
        "4 двигателя · 2 драйвера · камера · датчики",
        "4 motors · 2 drivers · camera · sensors",
    ),

    # Slide 5 — architecture
    ("04 · АРХИТЕКТУРА", "04 · ARCHITECTURE"),
    ("Как устроена система", "How the system is built"),
    ("Каждый модуль отвечает за свою задачу", "Each module has a clear role"),
    ("КАМЕРА", "CAMERA"),
    ("штрихкод · видеопоток", "barcode · video stream"),
    ("киоск · карта · маршрут", "kiosk · map · route"),
    ("привод · энкодеры", "drive · encoders"),
    ("ПРИВОД", "DRIVE"),
    ("4 мотора", "4 motors"),
    ("ПАССАЖИРСКИЙ ИНТЕРФЕЙС", "PASSENGER INTERFACE"),
    ("4 языка · билет · выбор пункта", "3 languages · ticket · destination choice"),
    (
        "Вычисления разделены: высокоуровневая логика — на Raspberry Pi,",
        "Computation is split: high-level logic runs on the Raspberry Pi,",
    ),
    (
        "реальное управление движением — на Arduino.",
        "real-time motion control runs on the Arduino.",
    ),

    # Slide 6 — navigation
    ("05 · НАВИГАЦИЯ", "05 · NAVIGATION"),
    ("От линии — к автономному маршруту", "From teleop to an autonomous route"),
    (
        "Что уже работает, что сделано в программе, что ещё нужно проверить",
        "What already works, what is implemented in software, what still needs proof",
    ),
    ("Работает в прототипе", "Works on the prototype"),
    ("Езда по стрелкам", "Arrow teleoperation"),
    ("4 стрелки ведут робота", "Four arrows drive the robot"),
    ("Реализовано в ПО", "Implemented in software"),
    ("Карта и маршрут", "Map and route"),
    (
        "Робот строит карту по LiDAR и ищет путь A→B.",
        "The robot builds a LiDAR map and plans path A→B.",
    ),
    ("Нужно подтвердить", "Still to confirm"),
    ("Полевые испытания", "Field trials"),
    (
        "Проверка в реальном зале с людьми.",
        "Validation in a real hall with people.",
    ),
    (
        "Программа навигации уже написана; работу в настоящем аэропорту пока не заявляем.",
        "Navigation software is already written; real airport operation is not claimed yet.",
    ),

    # Slide 7 — interface
    ("06 · ИНТЕРФЕЙС", "06 · INTERFACE"),
    ("Экран, понятный с первого взгляда", "A screen clear at first glance"),
    (
        "Крупные действия, четыре языка и быстрый доступ к билету",
        "Large actions, multiple languages, and quick ticket access",
    ),
    ("4 языка", "3 languages"),
    ("RU · EN · 中文 · TT", "RU · EN · TT"),
    ("6 направлений", "6 destinations"),
    ("4 из них настраиваются", "4 of them are configurable"),
    ("Билет", "Ticket"),
    ("камера или ручной ввод", "camera or manual entry"),

    # Slide 8 — safety
    ("07 · БЕЗОПАСНОСТЬ", "07 · SAFETY"),
    ("Безопасность заложена в сценарий", "Safety is built into the scenario"),
    (
        "Робот должен помогать, не создавая новых рисков",
        "The robot must help without creating new risks",
    ),
    ("ВНУТРЕННОСТИ ПРОТОТИПА", "PROTOTYPE INTERNALS"),
    ("Watchdog привода", "Drive watchdog"),
    ("Нет команд → остановка", "No commands → stop"),
    ("Heartbeat киоска", "Kiosk heartbeat"),
    (
        "Потеря связи → отмена сопровождения",
        "Link loss → cancel escort",
    ),
    ("Контроль маршрута", "Route control"),
    (
        "Только заранее настроенные точки",
        "Only preconfigured waypoints",
    ),
    ("Защитный корпус", "Protective body"),
    ("Нужно до эксплуатации", "Required before deployment"),
    ("Индикация", "Status lights"),
    ("Показывать состояние робота", "Show the robot state"),

    # Slide 9 — development
    ("08 · РАЗРАБОТКА", "08 · DEVELOPMENT"),
    ("Что было сделано в проекте", "What has been done"),
    ("Сценарий", "Scenario"),
    ("Определён путь пассажира", "Passenger journey defined"),
    ("Шасси", "Chassis"),
    ("Собрана Mecanum-платформа", "Mecanum platform assembled"),
    ("Электроника", "Electronics"),
    ("Подключены Pi, Mega и привод", "Pi, Mega, and drive connected"),
    ("Движение", "Motion"),
    ("Удалённое управление роботом", "Remote robot control"),
    ("Развитие", "Growth"),
    ("Добавлены киоск и LiDAR-модули", "Kiosk and LiDAR modules added"),
    (
        "Модульный подход позволяет развивать части проекта независимо",
        "A modular approach lets each part evolve independently",
    ),

    # Slide 10 — economics (use speech figure: ~90,000)
    ("09 · ЭКОНОМИКА", "09 · ECONOMICS"),
    (
        "Смета исходного прототипа — около 18 000 ₽",
        "Current prototype cost — about 90,000 ₽",
    ),
    ("75 000", "90 000"),
    ("рублей", "rubles"),
    (
        "Шасси · двигатели · драйверы · Arduino Mega",
        "Chassis · motors · drivers · Arduino Mega",
    ),
    (
        "Raspberry Pi 5 · камера · датчики · материалы",
        "Raspberry Pi 5 · camera · sensors · materials",
    ),
    ("Почему это важно", "Why it matters"),
    (
        "доступная база для школьного прототипа",
        "an affordable base for a school prototype",
    ),
    ("модульная замена компонентов", "modular component upgrades"),
    (
        "масштабирование по мере испытаний",
        "scaling as trials progress",
    ),

    # Slide 11 — results
    ("10 · РЕЗУЛЬТАТ", "10 · RESULTS"),
    (
        "Что уже достигнуто — и что дальше",
        "What is already achieved — and what is next",
    ),
    ("РЕАЛИЗОВАНО / ДЕМО", "DONE / DEMO"),
    ("манёвренная Mecanum-платформа", "maneuverable Mecanum platform"),
    ("разделённое управление Pi + Arduino", "split control: Pi + Arduino"),
    (
        "киоск и билетный сценарий на демоданных",
        "kiosk and ticket flow on demo data",
    ),
    (
        "watchdog, heartbeat и контроль точек",
        "watchdog, heartbeat, and waypoint control",
    ),
    ("ТРЕБУЕТ ПОДТВЕРЖДЕНИЯ", "NEEDS CONFIRMATION"),
    (
        "испытания LiDAR-навигации в терминале",
        "LiDAR navigation trials in a terminal",
    ),
    (
        "устойчивость к динамичным препятствиям",
        "robustness to dynamic obstacles",
    ),
    (
        "интеграция с реальными данными аэропорта",
        "integration with real airport data",
    ),
    ("промышленный безопасный корпус", "industrial safety body"),

    # Slide 12 — closing
    ("Технология,", "Technology"),
    ("которая помогает", "that helps people"),
    ("не потеряться", "not get lost"),
    (
        "Цель — передать роботу часть рутинной навигации",
        "The goal is to hand routine navigation to the robot",
    ),
    (
        "и потенциально снизить нагрузку на сотрудников.",
        "and potentially reduce staff workload.",
    ),
    ("Спасибо за внимание", "Thank you for your attention"),
]


def replace_in_paragraph(paragraph, mapping: dict[str, str]) -> bool:
    full = "".join(r.text for r in paragraph.runs) if paragraph.runs else (paragraph.text or "")
    if not full:
        return False
    new = mapping.get(full)
    if new is None:
        changed = False
        for old, new_val in mapping.items():
            if old in full and old != new_val:
                full = full.replace(old, new_val)
                changed = True
        if not changed:
            return False
        new = full

    # Keep two-tone brand split: Air + Porter (cyan on second run).
    if new == "AirPorter" and len(paragraph.runs) >= 2:
        paragraph.runs[0].text = "Air"
        paragraph.runs[1].text = "Porter"
        for r in paragraph.runs[2:]:
            r.text = ""
        return True

    if not paragraph.runs:
        paragraph.text = new
        return True
    paragraph.runs[0].text = new
    for r in paragraph.runs[1:]:
        r.text = ""
    return True


def build_pptx() -> Path:
    if OUT.exists():
        OUT.unlink()
    shutil.copy2(SRC, OUT)
    prs = Presentation(str(OUT))
    # longer keys first for partial replace safety
    mapping = dict(sorted(REPLACEMENTS, key=lambda kv: len(kv[0]), reverse=True))
    changed = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for p in shape.text_frame.paragraphs:
                if replace_in_paragraph(p, mapping):
                    changed += 1
    prs.save(str(OUT))
    print(f"PPTX saved: {OUT} ({changed} text blocks updated)")
    return OUT


POLISHED_EN_SPEECH = """AirPorter — presentation speech

Slide 1. Introduction to the project
AirPorter — a mobile robot assistant for the airport
Good afternoon! We present the AirPorter project — a mobile robot assistant for the airport.
The goal of the project is to make passengers’ time at the airport more comfortable. The robot helps people find their way around the terminal, escorts them to the required point, and can carry luggage — reducing physical strain and stress, especially for first-time visitors.

Slide 2. The problem
A modern airport is a complex system with a large volume of information, services, and zones. For many passengers — especially international travelers or those who fly rarely — finding their way becomes a serious challenge.
We identified three main difficulties:
• physical strain when carrying luggage;
• complex navigation inside the terminal;
• the language barrier when obtaining information.
These problems became the foundation of our project.

Slide 3. The solution
We propose a robot assistant that helps a passenger travel from a starting point to a destination according to their needs.
The system works in several stages.
Stage one — ticket scanning. The robot receives key information about the passenger and their flight.
Stage two — displaying information. The screen shows the flight number, boarding time, gate, and available services.
Then three operating scenarios are possible.
Scenario one. The passenger only needs information. The robot shows the route and the required data, and the session ends.
Scenario two. The passenger needs escorting. The robot suggests a route via selected points: check-in, restrooms, Duty Free, café, or the boarding gate. After the route is chosen, the robot escorts the passenger.
Scenario three. The passenger needs help only with luggage transport. The robot takes the bags and goes through all baggage-screening stages; the passenger only needs to complete personal security screening and board the aircraft.
In the second and third scenarios the robot can handle luggage. If full escorting is not required, it can also serve as a temporary place to store belongings.

Slide 4. Hardware platform
To implement the project, we built the first working prototype.
The main computing module is a Raspberry Pi 5 (16 GB). It runs ROS 2 Jazzy, processes lidar and camera data, and handles high-level logic.
Motor control is performed by an Arduino Mega 2560: it processes encoder signals and drives the platform.
We use L298N drivers, JGB37-520 motors, Mecanum wheels, a 10.1" touchscreen, a camera, and a COIN D6 lidar.

Slide 5. Navigation
We implemented indoor mapping with lidar. The robot estimates its position, accounts for its body dimensions and encoder data, builds a route, and avoids obstacles. It supports travel through waypoints and operator supervision.

Slide 6. Interface
The home screen presents the main airport services: check-in, boarding, baggage claim, the information desk, and other terminal points.
Russian, English, and Tatar are supported. The interface follows an airport visual theme and is designed for simple passenger use.

Slide 7. Safety
We implemented automatic stop on loss of connection, route monitoring, obstacle recognition, a protective acrylic body, low ground clearance, and status lighting.

Slide 8. First prototype
Today we present our team’s first prototype: an interactive on-screen menu, travel through waypoints, spatial orientation, and the ability to carry a load.

Slide 9. Expert feedback
The project was presented to aviation-industry experts and received positive feedback. Airport representatives expressed interest in further testing of the development.

Slide 10. Economics
The current prototype costs about 90,000 rubles. The modular design allows the system to be upgraded without fully replacing the hardware.

Slide 11. Future plans
Our plans are large and ambitious: substantial and exciting work on the project still lies ahead.
First, we will reinforce the structure with metal profiles. This will increase payload capacity and impact resistance.
Changes will also be needed in the wheelbase: we plan industrial metal Mecanum wheels.
It is important to keep safety and a neat appearance: with the move to metal we will add more protective and cosmetic body details.
On the electronics side, we plan a gradual shift to industrial-grade components to reach real airport trials sooner.
In navigation, one of the key goals is stronger sensing, including a 360° industrial camera and additional modules.
We will also develop convenient interfaces for passengers and for the development team.
For users — an app showing where the robot is, what stage of service it is in, and baggage alerts when needed.
For developers — a fleet dashboard with system health, occupancy, and possible faults.
This will let the team pinpoint issues, fix them remotely, or send a robot to a specialized hub for repair.
The hub will become a charging station and a holding area for robots awaiting service, from which they can be taken to the workshop for repair and cleaning.
Separately, we plan heating and cooling for drinks in the existing cup holder.
The final stage is testing the system in real airport scenarios.

Slide 13. Closing
Thank you for your attention!
The AirPorter Team
"""


def update_docx() -> Path:
    """Rewrite the English half of the speech DOCX from polished text."""
    try:
        from docx import Document
    except ImportError:
        print("python-docx not installed; writing polished EN speech as .txt")
        txt = DOCX_OUT.with_suffix(".en.txt")
        txt.write_text(POLISHED_EN_SPEECH, encoding="utf-8")
        return txt

    shutil.copy2(DOCX_SRC, DOCX_OUT)
    doc = Document(str(DOCX_OUT))

    # Find start of English section
    en_start = None
    for i, p in enumerate(doc.paragraphs):
        if p.text.strip().startswith("AirPorter — presentation speech"):
            en_start = i
            break
    if en_start is None:
        raise RuntimeError("English section not found in DOCX")

    en_lines = [ln for ln in POLISHED_EN_SPEECH.splitlines() if ln.strip()]
    # Clear English paragraphs and refill sequentially
    en_paras = doc.paragraphs[en_start:]
    for p in en_paras:
        for run in p.runs:
            run.text = ""
        if p.text:
            p.text = ""

    # Put lines into existing paragraph slots; append if needed
    for idx, line in enumerate(en_lines):
        if idx < len(en_paras):
            p = en_paras[idx]
            if p.runs:
                p.runs[0].text = line
            else:
                p.add_run(line)
        else:
            doc.add_paragraph(line)

    doc.save(str(DOCX_OUT))
    print(f"DOCX saved: {DOCX_OUT}")
    return DOCX_OUT


def main() -> None:
    build_pptx()
    update_docx()


if __name__ == "__main__":
    main()
