# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Emu
from pathlib import Path

pptx = Path(r"c:\Users\user\Downloads\Telegram Desktop\AeroPorter_презентация (2).pptx")
prs = Presentation(str(pptx))
print("slides", len(prs.slides))
print("size", prs.slide_width, prs.slide_height)
print("w_in", prs.slide_width / 914400, "h_in", prs.slide_height / 914400)

for i, slide in enumerate(prs.slides, 1):
    print("=" * 60)
    print("SLIDE", i, "layout", slide.slide_layout.name)
    for shape in slide.shapes:
        bits = [
            f"name={shape.name!r}",
            f"type={shape.shape_type}",
            f"L={shape.left}",
            f"T={shape.top}",
            f"W={shape.width}",
            f"H={shape.height}",
        ]
        if hasattr(shape, "image"):
            try:
                bits.append(f"img={shape.image.content_type} {len(shape.image.blob)}b")
            except Exception:
                pass
        if shape.has_text_frame:
            texts = []
            for p in shape.text_frame.paragraphs:
                t = "".join(r.text for r in p.runs) if p.runs else (p.text or "")
                if t.strip():
                    texts.append(t.strip())
            if texts:
                bits.append("TEXT=" + " || ".join(texts)[:240])
        print(" ", " | ".join(bits))
